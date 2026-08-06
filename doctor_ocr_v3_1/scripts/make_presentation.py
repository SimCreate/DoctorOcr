#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
DoctorOcr 발표 PPT 생성기 (10분, 11장)
- 인물: 야간대학원생 / 기존 코드 재현 + 검증 + 개선 스토리
- 핵심: 98.8% 리키지 폭로 → 클린 재실험 → attention+증강 재도전
- CTC는 가볍게, v3_1(진행 중)이 스토리의 중심
실행: doctor_ocr_v3_1/venv/bin/python scripts/make_presentation.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path("/home/dev/DoctorOcr/doctor_ocr_v3_1/DoctorOcr_발표_20260807.pptx")

# ---- 팔레트 ----
DARK   = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT = RGBColor(0x2E, 0x86, 0xDE)   # 파랑 (attention/개선)
RED    = RGBColor(0xC0, 0x39, 0x2B)   # 리키지/실패
GREEN  = RGBColor(0x27, 0xAE, 0x60)   # 성공
GRAY   = RGBColor(0x6C, 0x75, 0x7D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF0, 0xF3, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=LIGHT):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def txt(slide, l, t, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="맑은 고딕"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def bullets(slide, l, t, w, h, items, size=16, color=DARK, gap=6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (mark, text, c, b) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r1 = p.add_run()
        r1.text = mark + " " if mark else ""
        r1.font.size = Pt(size)
        r1.font.bold = (mark in ("▣", "▶"))
        r1.font.color.rgb = c
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(size)
        r2.font.bold = b
        r2.font.color.rgb = color
        r2.font.name = "맑은 고딕"
    return tb


def header(slide, title, sub=None, accent=ACCENT):
    bg(slide)
    box(slide, 0, 0, prs.slide_width, Inches(1.0), fill=accent)
    txt(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.8),
        title, size=26, color=WHITE, bold=True)
    if sub:
        txt(slide, Inches(0.5), Inches(0.78), Inches(12), Inches(0.3),
            sub, size=12, color=RGBColor(0xEF, 0xF1, 0xF6))


def footer(slide, n):
    txt(slide, prs.slide_width - Inches(1.2), Inches(7.05), Inches(1.0), Inches(0.3),
        f"{n}", size=11, color=GRAY, align=PP_ALIGN.RIGHT)


def stat_box(slide, l, t, w, h, num, label, color=ACCENT):
    box(slide, l, t, w, h, fill=color)
    txt(slide, l + Inches(0.1), t + Inches(0.08), w - Inches(0.2), Inches(0.5),
        num, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, l + Inches(0.05), t + Inches(0.55), w - Inches(0.1), Inches(0.4),
        label, size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ============ 1. 표지 ============
s = add_slide(); bg(s, DARK)
box(s, 0, Inches(4.6), prs.slide_width, Inches(2.9), fill=RGBColor(0x27, 0x2A, 0x3D))
txt(s, Inches(0.8), Inches(2.0), Inches(12), Inches(1.2),
    "의사 손글씨 처방전 OCR", size=44, color=WHITE, bold=True)
txt(s, Inches(0.8), Inches(3.1), Inches(12), Inches(0.8),
    "재현 · 검증 · 개선 — 98.8%의 진실과 재도전", size=24, color=ACCENT, bold=True)
txt(s, Inches(0.8), Inches(5.1), Inches(12), Inches(0.6),
    "야간대학원 연구진행 발표  |  2026. 08. 07", size=18, color=WHITE)
txt(s, Inches(0.8), Inches(5.8), Inches(12), Inches(0.5),
    "RxHandBD (Kaggle) 5,578장  ·  CRNN 계열  ·  v1 → v2.2 → v3 계보", size=13, color=RGBColor(0x9A, 0xA0, 0xB0))
footer(s, 1)

# ============ 2. 배경 & 문제 ============
s = add_slide(); header(s, "문제 정의 — 의사 필기 처방전은 왜 어려운가")
bullets(s, Inches(0.6), Inches(1.4), Inches(12), Inches(4.5), [
    ("", "RxHandBD: 실제 의사 처방전 스캔 5,578장, 라벨은 1,788개 고유 단어", DARK, False),
    ("", "▶ 롱테일 분포: 고유 라벨의 64%가 '딱 1번' 등장 → 모델이 거의 못 본 단어가 절반 이상", ACCENT, True),
    ("", "키워드: 의약품명 (Napa, Devixil, Ecosprin Plus …) — 인식은 커녕 철자 하나 틀려도 위험", DARK, False),
    ("", "인식 정확도 지표: exact match (단어 전체 일치) + CER (문자 오류율)", DARK, False),
    ("", "이미지 단위 평가가 아니라 '단어' 단위 평가가 실제 의료 활용에 맞음", DARK, False),
])
stat_box(s, Inches(0.7), Inches(5.9), Inches(3.7), Inches(1.1), "5,578장", "데이터", red if False else RGBColor(0x24, 0x55, 0x7E))
stat_box(s, Inches(4.7), Inches(5.9), Inches(3.7), Inches(1.1), "1,788", "고유 라벨", ACCENT)
stat_box(s, Inches(8.7), Inches(5.9), Inches(3.7), Inches(1.1), "64%", "1회 등장 (롱테일)", RED)
footer(s, 2)

# ============ 3. 왜 CRNN 계열인가 ============
s = add_slide(); header(s, "Architecture — CRNN: 이미지를 문자 시퀀스로")
bullets(s, Inches(0.6), Inches(1.3), Inches(7.2), Inches(5.0), [
    ("O", "CNN 인코더: 이미지 → 특징맵 (어디에 뭐가 그려졌는가)", DARK, False),
    ("O", "BiLSTM: 특징을 '왼쪽→오른쪽, 오른쪽→왼쪽' 시퀀스로", DARK, False),
    ("O", "디코더가 문자열을 만들며 시퀀스를 단어로 변환", DARK, False),
    ("", "CRNN(Shi et al. 2015)은 손글씨/문자인식의 고전적 표준 구조", GRAY, False),
    ("", "→ 왜 순서(문자열)를 다루는 모델이 필요한가: 'Nap' vs 'Npa'", GRAY, False),
])
# 단순 구조 다이어그램 (텍스트 박스 상자)
box(s, Inches(8.4), Inches(1.6), Inches(1.9), Inches(1.2), fill=RGBColor(0xD5, 0xE8, 0xD4))
txt(s, Inches(8.4), Inches(1.85), Inches(1.9), Inches(0.8), "CNN\n인코더", size=13, bold=True, align=PP_ALIGN.CENTER)
box(s, Inches(10.6), Inches(1.6), Inches(2.0), Inches(1.2), fill=RGBColor(0xFF, 0xE9, 0xC7))
txt(s, Inches(10.6), Inches(1.85), Inches(2.0), Inches(0.8), "BiLSTM\n시퀀스", size=13, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(9.55), Inches(1.35), Inches(0.5), Inches(0.4), "→", size=20, align=PP_ALIGN.CENTER)
txt(s, Inches(11.8), Inches(1.35), Inches(0.5), Inches(0.4), "→", size=20, align=PP_ALIGN.CENTER)
box(s, Inches(8.4), Inches(3.1), Inches(4.2), Inches(1.0), fill=RGBColor(0xD6, 0xE4, 0xFF))
txt(s, Inches(8.4), Inches(3.3), Inches(4.2), Inches(0.7), "디코더 → 문자열\n(CTC or Attention)", size=12, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(8.4), Inches(4.3), Inches(4.3), Inches(0.4), "디코더 방식이 이 연구의 핵심 축", size=11, color=RED, align=PP_ALIGN.CENTER)
footer(s, 3)

# ============ 4. 기존 코드 재현 & '98.8%' 의심 ============
s = add_slide(); header(s, "재현 — 기존 파이프라인을 그대로 돌렸더니 '98.8%' ?", accent=RED)
bullets(s, Inches(0.6), Inches(1.4), Inches(12), Inches(3.9), [
    ("", "v1 → v2 → v2.2 계보의 CRNN 코드를 로컬 GPU에서 재현", DARK, False),
    ("", "▶ v2.2(CTC) + 증강2배 실험이 '원본 val에서 정확도 98.8%, PASS'로 기록됨", RED, True),
    ("", "한편 클린 기준(리키지 제거) 재측정에서는 전부 FAIL (35.8% / 37.9% / 36.4%)", RED, False),
    ("", "", GRAY, False),
], )
box(s, Inches(0.7), Inches(4.6), Inches(12.0), Inches(2.3), fill=RGBColor(0xFD, 0xEC, 0xEA))
txt(s, Inches(1.0), Inches(4.8), Inches(11.4), Inches(0.5), "핵심 의심 — 이미지 단위 리키지", size=18, color=RED, bold=True)
bullets(s, Inches(1.0), Inches(5.3), Inches(11.4), Inches(1.4), [
    ("", "기존 split은 '증강 포함 전체'를 80/20으로 나눔 → val 1,116장 중 약 890장(80%)이 train에 포함", DARK, False),
    ("", "→ '98.8%'는 검증 이미지를 이미 학습한 모델을 측정한 수치", RED, True),
])
footer(s, 4)

# ============ 5. 리키지 검증 & 클린 재실험 ============
s = add_slide(); header(s, "검증 — 리키지를 차단하고 다시 측정했다", accent=GREEN)
bullets(s, Inches(0.6), Inches(1.4), Inches(12), Inches(4.4), [
    ("", "클린 스플릿 설계: 원본을 80/20으로 먼저 분리 → train에만 증강/합성", DARK, False),
    ("", "val 1,116장이 어떤 실험군의 train에도 1장도 포함되지 않음을 스크립트로 대조", DARK, False),
    ("", "리키지 가드: train∩val > 0 이면 학습 중단하도록 파이프라인 보강", DARK, False),
    ("", "▶ 결과: 전 실험군 FAIL — 증강 효과는 +2.2p에 불과 (35.8→37.9%)", ACCENT, True),
])
import csv
res = {
    "exp1 원본": ("35.8%", "FAIL"), "exp2 증강": ("37.9%", "FAIL"), "exp3 +합성": ("36.4%", "FAIL")
}
x = Inches(0.7)
for k, (acc, verdict) in res.items():
    box(s, x, Inches(5.9), Inches(3.7), Inches(1.0), fill=RGBColor(0x2E, 0x40, 0x55))
    txt(s, x + Inches(0.2), Inches(6.0), Inches(2.2), Inches(0.4), k, size=13, color=WHITE)
    txt(s, x + Inches(2.0), Inches(6.0), Inches(1.5), Inches(0.4), acc, size=20, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    txt(s, x + Inches(0.2), Inches(6.45), Inches(3.3), Inches(0.4), verdict, size=14, color=RGBColor(0xFF, 0x8A, 0x80), bold=True)
    x += Inches(4.1)
footer(s, 5)

# ============ 6. 왜 CTC를 버리려 하는가 ============
s = add_slide(); header(s, "회고 — '가장 진보' vs '가장 잘 도는' 디코더", accent=ACCENT)
bullets(s, Inches(0.6), Inches(1.4), Inches(12), Inches(4.4), [
    ("", "v2.2(CTC): 정렬(alignment) 학습 불필요 → 견고하지만, 표현력이 단순(1층 linear)", DARK, False),
    ("", "v2.1(Attention): multi-head 8 + beam search — 더 표현력 있으나 데이터가 많아야 학습됨", DARK, False),
    ("", "모델 계보를 돌아보면: 이론적으론 attention이 더 진보, 실제론 CTC가 먼저 '작동'했던 것 뿐", GRAY, False),
    ("", "▶ v3 실험이 증명: CTC는 데이터를 늘려도 37.9%에서 정체 → 모델(디코더) 축의 한계", ACCENT, True),
], )
stat_box(s, Inches(0.8), Inches(5.9), Inches(3.7), Inches(1.1), "CTC 37.9%", "v3 최선 (클린)", RGBColor(0x24, 0x55, 0x7E))
stat_box(s, Inches(4.8), Inches(5.9), Inches(3.7), Inches(1.1), "Attention?", "데이터가 있다면", ACCENT)
stat_box(s, Inches(8.8), Inches(5.9), Inches(3.7), Inches(1.1), "가설 검증", "v3_1 실험", GREEN)
footer(s, 6)

# ============ 7. 개선 시도 — v3_1: attention + 증강 데이터 ============
s = add_slide(); header(s, "가설 실험 v3_1 — '진보한 디코더 + 충분한 데이터'", accent=ACCENT)
bullets(s, Inches(0.6), Inches(1.4), Inches(12), Inches(4.6), [
    ("", "가설: attention 디코더 실패의 원인은 '데이터 부족' — 데이터를 늘리면 수렴할 것", ACCENT, True),
    ("", "실험 구성 (기존의 실패 요인을 하나씩 제거):", DARK, False),
    ("", "  ① 데이터: v3 증강2배 (원본 4,462 + 증강 8,924 = 13,386장) — train만", DARK, False),
    ("", "  ② val: 클린 고정 split (1,116장) — v3와 동일 기준, 직접 비교 가능", DARK, False),
    ("", "  ③ 런타임 증강 OFF — 온디스크 증강과 중복 방지", DARK, False),
    ("", "  ④ checkpoint config 버그 수정 (hidden 384 실제값 저장)", DARK, False),
], )
box(s, Inches(0.7), Inches(6.0), Inches(12.0), Inches(0.9), fill=RGBColor(0xE, 0x1, 0x45, 0x5E) if False else RGBColor(0xDD, 0xE8, 0xF7))
txt(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.5),
    "실행 환경: GPU1 (Max-Q)  ·  batch 40  ·  effective 160  ·  warmup+cosine  ·  teacher forcing 0.47→0", size=13, color=RGBColor(0x1B, 0x3A, 0x5C))
footer(s, 7)

# ============ 8. 실험 결과 ============
s = add_slide(); header(s, "실험 결과 — attention은 CTC를 넘지 못했다", accent=ACCENT)
bullets(s, Inches(0.6), Inches(1.3), Inches(12), Inches(3.5), [
    ("", "v3_1(attention+증강)도 클린 val 평가에서 FAIL — 수용기준 미충족", RED, True),
    ("", "exact 30.4% / CER 38.2% — 학습 로그(argmax 43%)보다 실제 디코딩이 낮음", DARK, False),
    ("", "고빈도 54.8%는 CTC(56.5%)와 비슷, 그러나 중빈도 3.2%·저빈도 0.4%로 붕괴", RED, True),
    ("", "→ 많은 단어는 그럭저럭, 드문 단어는 통째로 오생성 (attention 특성)", GRAY, False),
])
# 비교 박스 3개
box_c = [
    ("v3 exp2 (CTC)", "37.9%", "CER 20.6%", "중빈도 33.5%", RGBColor(0x27,0xAE,0x60)),
    ("v3_1 (Attention)", "30.4%", "CER 38.2%", "중빈도  3.2%", RGBColor(0xC0,0x39,0x2B)),
]
x = Inches(1.0)
for name, exact, cer, mid, col in box_c:
    box(s, x, Inches(4.7), Inches(5.3), Inches(1.6), fill=RGBColor(0x2E,0x40,0x55))
    txt(s, x+Inches(0.2), Inches(4.8), Inches(4.9), Inches(0.4), name, size=16, color=WHITE, bold=True)
    txt(s, x+Inches(0.2), Inches(5.2), Inches(4.9), Inches(0.45), f"exact {exact}  {cer}", size=18, color=col, bold=True)
    txt(s, x+Inches(0.2), Inches(5.7), Inches(4.9), Inches(0.4), mid, size=13, color=RGBColor(0xCE,0xD4,0xE0))
    x += Inches(6.3)
footer(s, 8)

# ============ 9. 실험의 의미 ============
s = add_slide(); header(s, "실험의 의미 — 무엇이 검증됐나", accent=ACCENT)
bullets(s, Inches(0.6), Inches(1.4), Inches(12), Inches(4.4), [
    ("", "긍정: attention 디코더는 데이터가 커지면 수렴함 — v2.1 원본(20%) 대비 30.4%로 상승", GREEN, True),
    ("", "그러나: CTC가 여전히 실질 우위 — exact 37.9%, CER 20.6% vs attention 30.4%", ACCENT, True),
    ("", "핵심 인사이트: 디코더 선택보다, '어떤 오류가 더 치명적인가'가 중요", DARK, False),
    ("", "  CTC는 중·저빈도도 부분 일치(중 33.5%) → CER 낮음, attention은 저빈도 전멸(0.4%) → CER 급증", DARK, False),
    ("", "결론: '데이터가 attention을 살린다'는 가설은 부분 검증, 'attention이 CTC보다 좋다'는 기각", RED, True),
])
footer(s, 9)

# ============ 10. 결론 & 배운 것 ============
s = add_slide(); header(s, "결론 — 재현과 검증이 만든 교훈", accent=DARK)
bullets(s, Inches(0.6), Inches(1.4), Inches(12), Inches(4.4), [
    ("1", "검증 우선: 98.8%는 리키지 착시 — '정확도 숫자'보다 split 무결성이 먼저", RED, True),
    ("2", "데이터와 모델의 관계: 증강은 CTC에서 +2.2p, attention은 학습을 살렸으나 최종 성능은 CTC 미달 — 디코더 선택이 결정적", ACCENT, True),
    ("3", "롱테일(1회 등장 64%)은 데이터량으로만 안 풀림 → 라벨 정제/도메인 지식이 다음 축", GRAY, True),
    ("4", "야간 대학원 현실: 시간이 한정 → '검증 가능한 실험' 단위로 진행, 발표는 그 기록", GRAY, False),
])
footer(s, 10)

# ============ 11. Q&A ============
s = add_slide(); bg(s, DARK)
txt(s, Inches(0.8), Inches(3.0), Inches(12), Inches(1.2), "Q&A", size=48, color=WHITE, bold=True)
txt(s, Inches(0.8), Inches(4.2), Inches(12), Inches(0.6),
    "감사합니다", size=20, color=ACCENT)
footer(s, 11)

prs.save(OUT)
print(f"저장: {OUT}")
